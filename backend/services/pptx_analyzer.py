from pptx import Presentation


class PPTXAnalyzer:

    def extract(self, file_path: str):

        presentation = Presentation(file_path)

        content = []

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    content.append(shape.text)

        return "\n".join(content)